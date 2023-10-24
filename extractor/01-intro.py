import os
from pathlib import Path
from pptx import Presentation

import os
script_directory = os.path.dirname(os.path.abspath(__file__))
path = Path(script_directory).parent / 'pptx' / '02-Physical_Spiritual.pptx'

prs = Presentation(path)

# Initialize variables to keep track of sequence and slide number
sequence = 1
slide_number = 1

# Initialize an empty list to store the HTML/Markdown content
content = []

# Iterate through slides in the presentation
for slide in prs.slides:
    for shape in slide.shapes:
        # Check if the shape is a text box
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if text_frame.paragraphs:
                # Extract text from the text box
                text = ''
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text

                # Extract hyperlinks (if any)
                hyperlinks = []
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink._hlinkClick:
                            for hyperlink in run.hyperlink._hlinkClick:
                                hyperlinks.append(hyperlink.address)

                # Generate HTML/Markdown for the text box
                content.append(f"<div v-click='{sequence}' class='text-xs group/ii'>{text}")
                for i, hyperlink in enumerate(hyperlinks):
                    content.append(f"<a href='{hyperlink}' v-click='{sequence + i + 1}'></a>")
                content.append("</div>")

                # Increment the sequence
                sequence += len(hyperlinks) + 1

    # Increment slide number
    slide_number += 1

# Convert the content list to a single string
html_output = '\n'.join(content)

# Print or save the HTML/Markdown output
print(html_output)

# You can write the HTML output to a file if needed
# with open('output.html', 'w', encoding='utf-8') as f:
#     f.write(html_output)
