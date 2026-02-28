#!/usr/bin/env python3
# Licensed under MIT License. See LICENSE file for details.
# https://github.com/matttrice/hsu-extractor

import os
import glob
import html
import json
import math
import re
import shutil
import sys
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

# XML namespaces used in PPTX files
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# Target canvas dimensions for MBS (960×540 pixels, 16:9 aspect ratio)
TARGET_CANVAS_WIDTH = 960
TARGET_CANVAS_HEIGHT = 540

# EMU to pixels conversion (96 DPI standard)
EMU_PER_PIXEL = 9525

# Wrap inference heuristic defaults
WRAP_MIN_TEXT_LENGTH = 40
WRAP_AVG_CHAR_WIDTH_EM = 0.52
WRAP_FALLBACK_FONT_SIZE_PX = 20

SCRIPTURE_REFERENCE_PLUS_TEXT_RE = re.compile(
    r"^\s*(?:[1-3]\s+)?[A-Za-z][A-Za-z'’.-]*(?:\s+[A-Za-z][A-Za-z'’.-]*)*\s+"
    r"\d{1,3}:\d{1,3}(?:[-–]\d{1,3})?(?:\s*,\s*\d{1,3}(?:[-–]\d{1,3})?)*"
    r"(?:\s*[:—-]\s*|\s+).+",
    re.IGNORECASE,
)


def normalize_unicode_text(text):
    """Replace common Unicode typographic characters with plain ASCII equivalents."""
    if not text:
        return text
    # Smart double quotes → straight double
    text = text.replace('\u201c', '"').replace('\u201d', '"')
    # Smart single quotes / apostrophes → straight single
    text = text.replace('\u2018', "'").replace('\u2019', "'")
    # En-dash / em-dash → hyphen / double-hyphen
    text = text.replace('\u2013', '-').replace('\u2014', '--')
    # Ellipsis → three dots
    text = text.replace('\u2026', '...')
    # Non-breaking space → regular space
    text = text.replace('\u00a0', ' ')
    return text


def _looks_like_scripture_reference_with_text(text):
    """Best-effort check for leading scripture reference followed by content.

    Matches patterns like:
    - "Ezekiel 18:4: For every living soul..."
    - "Genesis 27:30-38: After Isaac finished..."
    """
    if not isinstance(text, str) or not text.strip():
        return False

    single_line = ' '.join(text.split())
    return bool(SCRIPTURE_REFERENCE_PLUS_TEXT_RE.match(single_line))


def _sanitize_scripture_block_font(entry, is_scripture=False):
    """Remove block-level typography that should be delegated to inline scripture markup.

    Scripture text often contains mixed run formatting (e.g., bold/large reference + normal body).
    Setting a single block font_size/bold at shape level causes incorrect rendering downstream.

    Also propagates `is_scripture` flag to the JSON entry so downstream consumers
    can identify scripture blocks structurally rather than by parsing rendered markup.

        Contract:
        - This function must be driven by explicit scripture metadata (`is_scripture`),
            not by checking rendered markup.
        - Keep layout/flow font fields (e.g., wrap, align, v_align) intact.
        - Remove only typography fields that can incorrectly style an entire mixed-format block.
    """
    if not entry or not is_scripture:
        return

    # Hyperlinked entries are drill links to scripture routes, not scripture
    # content themselves — preserve their font and skip the flag.
    if entry.get('hyperlink'):
        return

    entry['is_scripture'] = True

    font = entry.get('font')
    if not isinstance(font, dict):
        return

    font.pop('font_size', None)
    font.pop('bold', None)

    if not font:
        entry.pop('font', None)

def extract_theme_colors_from_pptx(pptx_path):
    """Extract theme color scheme from PowerPoint file.
    
    Returns dict mapping MSO_THEME_COLOR string representations to hex colors.
    Example: "TEXT_1 (13)" -> "#000000"
    """
    theme_colors = {}
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # Read theme XML
            theme_xml = zip_ref.read('ppt/theme/theme1.xml')
            root = ET.fromstring(theme_xml)
            
            # Find color scheme
            clr_scheme = root.find('.//a:clrScheme', NAMESPACES)
            if clr_scheme is None:
                return {}
            
            # Map theme element names to MSO_THEME_COLOR enum indices
            # Based on Office Open XML spec and python-pptx MSO_THEME_COLOR enum
            color_mapping = [
                ('dk1', 'TEXT_1', 13),           # Dark 1 (main text)
                ('lt1', 'BACKGROUND_1', 14),     # Light 1 (main background)
                ('dk2', 'TEXT_2', 15),           # Dark 2 (secondary text)
                ('lt2', 'BACKGROUND_2', 16),     # Light 2 (secondary background)
                ('accent1', 'ACCENT_1', 5),      # Accent 1
                ('accent2', 'ACCENT_2', 6),      # Accent 2
                ('accent3', 'ACCENT_3', 7),      # Accent 3
                ('accent4', 'ACCENT_4', 8),      # Accent 4
                ('accent5', 'ACCENT_5', 9),      # Accent 5
                ('accent6', 'ACCENT_6', 10),     # Accent 6
                ('hlink', 'HYPERLINK', 11),      # Hyperlink
                ('folHlink', 'FOLLOWED_HYPERLINK', 12),  # Followed Hyperlink
            ]
            
            for elem_name, mso_name, mso_index in color_mapping:
                color_elem = clr_scheme.find(f'a:{elem_name}', NAMESPACES)
                if color_elem is not None:
                    # Try RGB color first
                    srgb = color_elem.find('.//a:srgbClr', NAMESPACES)
                    if srgb is not None:
                        hex_val = '#' + srgb.get('val', '').upper()
                        # Store with format that matches python-pptx output: "NAME (index)"
                        theme_colors[f"{mso_name} ({mso_index})"] = hex_val
                    else:
                        # Try system color
                        sys_clr = color_elem.find('.//a:sysClr', NAMESPACES)
                        if sys_clr is not None:
                            last_clr = sys_clr.get('lastClr')
                            if last_clr:
                                hex_val = '#' + last_clr.upper()
                                theme_colors[f"{mso_name} ({mso_index})"] = hex_val
    
    except Exception as e:
        print(f"Warning: Could not extract theme colors: {e}")
    
    return theme_colors


def resolve_theme_color(theme_color_str, theme_map, context=""):
    """Resolve a theme color string to hex value.
    
    Args:
        theme_color_str: Theme color string like "TEXT_1 (13)"
        theme_map: Dictionary mapping theme colors to hex values
        context: Context for warning message (e.g., "fill", "line", "font")
    
    Returns:
        Hex color string, or "#000000" if not found
    """
    if theme_color_str in theme_map:
        return theme_map[theme_color_str]
    
    # Theme color not in map - warn and default to black
    print(f"Warning: Unknown theme color '{theme_color_str}' in {context}. Using #000000 (black).")
    return "#000000"

def emu_to_px(emu):
    """Convert EMU to pixels at 96 DPI."""
    if emu is None:
        return None
    return round(emu / EMU_PER_PIXEL)

def scale_to_target(value, source_width, target_width=TARGET_CANVAS_WIDTH):
    """Scale a coordinate value from source slide dimensions to target canvas dimensions.
    
    Args:
        value: The coordinate value to scale (can be None)
        source_width: The width of the source PowerPoint slide in pixels
        target_width: The target canvas width (default: 960)
    
    Returns:
        Scaled value rounded to whole number, or None if value is None
    """
    if value is None:
        return None
    if source_width == target_width:
        return round(value)
    scale_factor = target_width / source_width
    return round(value * scale_factor)


def should_force_font_wrap(text, layout=None, font_size=None):
    """Infer whether text should be wrapped even when PPT wrap is not explicitly set.

    Heuristics used:
    1) Explicit line-break markers in text imply wrapping intent.
    2) Long single-line text that exceeds estimated line capacity implies wrap.
    """
    if not text:
        return False

    normalized_text = text.replace('\r\n', '\n').replace('\r', '\n').replace('\v', '\n')

    # Manual line breaks strongly indicate wrapped content intent.
    if '\n' in normalized_text:
        return True

    if not layout:
        return False

    width = layout.get('width')
    if not width or width <= 0:
        return False

    # Avoid forcing wrap for very short text.
    compact_text = ' '.join(normalized_text.split())
    if len(compact_text) < WRAP_MIN_TEXT_LENGTH:
        return False

    # Approximate average glyph width for mixed-case text.
    # WRAP_AVG_CHAR_WIDTH_EM is a conservative middle-ground for common sans-serif fonts.
    size_px = font_size if isinstance(font_size, (int, float)) and font_size > 0 else WRAP_FALLBACK_FONT_SIZE_PX
    avg_char_width = max(size_px * WRAP_AVG_CHAR_WIDTH_EM, 1)
    estimated_chars_per_line = max(int(width / avg_char_width), 1)

    # If a single visual line clearly exceeds the box capacity, force wrap.
    longest_line_len = max((len(part) for part in compact_text.split('\n')), default=0)
    return longest_line_len > estimated_chars_per_line

def rgb_to_hex(rgb_color):
    """Convert RGBColor to hex string."""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color}"
    except:
        return None

def enumerate_shapes_recursive(shapes, z_index_start=0, parent_group_id=None):
    """Recursively enumerate all shapes including those inside groups.
    
    Args:
        shapes: Collection of shapes to enumerate (from slide.shapes or group.shapes)
        z_index_start: Starting z-index for enumeration
        parent_group_id: Parent group's shape ID if applicable
    
    Yields:
        Tuple of (z_index, shape, group_id) for each shape found
    """
    z_idx = z_index_start
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # This is a group - recursively enumerate its children
            group_id = str(shape.shape_id)
            yield from enumerate_shapes_recursive(shape.shapes, z_idx, group_id)
            # Count group members for z-index offset
            z_idx += len(list(shape.shapes))
        else:
            # Regular shape
            yield (z_idx, shape, parent_group_id)
            z_idx += 1

def extract_layout_from_xml(shape_elem, slide_width=None):
    """Extract layout from shape XML element.
    
    Args:
        shape_elem: XML element (p:sp or p:cxnSp)
        slide_width: Source slide width for auto-scaling
    
    Returns:
        Layout dict with coordinates scaled to target canvas
    """
    try:
        spPr = shape_elem.find('p:spPr', NAMESPACES)
        if spPr is None:
            return None
        
        xfrm = spPr.find('a:xfrm', NAMESPACES)
        if xfrm is None:
            return None
        
        # Get offset (position)
        off = xfrm.find('a:off', NAMESPACES)
        # Get extents (size)
        ext = xfrm.find('a:ext', NAMESPACES)
        
        if off is None or ext is None:
            return None
        
        # Extract EMU values and convert to pixels
        x = emu_to_px(int(off.get('x', 0)))
        y = emu_to_px(int(off.get('y', 0)))
        width = emu_to_px(int(ext.get('cx', 0)))
        height = emu_to_px(int(ext.get('cy', 0)))
        rotation = int(xfrm.get('rot', 0)) / 60000  # Convert from 1/60000 degrees
        
        # Scale to target canvas
        if slide_width is not None:
            x = scale_to_target(x, slide_width)
            y = scale_to_target(y, slide_width)
            width = scale_to_target(width, slide_width)
            height = scale_to_target(height, slide_width)
        
        return {
            'x': x,
            'y': y,
            'width': width,
            'height': height,
            'rotation': rotation
        }
    except:
        return None

def extract_font_from_xml(shape_elem, layout=None):
    """Extract font properties from shape XML element.
    
    Args:
        shape_elem: XML element (p:sp)
    
    Returns:
        Font dict with properties
    """
    try:
        txBody = shape_elem.find('.//p:txBody', NAMESPACES)
        if txBody is None:
            return None
        
        font_data = {}
        
        # Get bodyPr for vertical alignment
        bodyPr = txBody.find('a:bodyPr', NAMESPACES)
        if bodyPr is not None:
            anchor = bodyPr.get('anchor')
            if anchor == 't':
                font_data['v_align'] = 'top'
            elif anchor == 'ctr':
                font_data['v_align'] = 'middle'
            elif anchor == 'b':
                font_data['v_align'] = 'bottom'

            # PowerPoint wrap settings:
            # - wrap="none" => no wrapping
            # - wrap missing or other values => wrapping enabled/default
            wrap_attr = bodyPr.get('wrap')
            if wrap_attr != 'none':
                font_data['wrap'] = True
        
        # Get first paragraph for alignment and default properties
        para = txBody.find('a:p', NAMESPACES)
        if para is not None:
            pPr = para.find('a:pPr', NAMESPACES)
            if pPr is not None:
                algn = pPr.get('algn')
                if algn == 'l':
                    font_data['align'] = 'left'
                elif algn == 'ctr':
                    font_data['align'] = 'center'
                elif algn == 'r':
                    font_data['align'] = 'right'
                
                # Check default run properties (defRPr) for inherited formatting
                defRPr = pPr.find('a:defRPr', NAMESPACES)
                if defRPr is not None:
                    # Font size from defRPr
                    sz = defRPr.get('sz')
                    if sz:
                        points = int(sz) / 100
                        font_size = round(points * 1.333, 1)
                        font_data['font_size'] = font_size
                    
                    # Bold from defRPr
                    if defRPr.get('b') == '1':
                        font_data['bold'] = True
                    
                    # Italic from defRPr
                    if defRPr.get('i') == '1':
                        font_data['italic'] = True

                    # Underline from defRPr
                    u_attr = defRPr.get('u')
                    if u_attr and u_attr != 'none':
                        font_data['underline'] = True
                    
                    # Font name from defRPr
                    latin = defRPr.find('a:latin', NAMESPACES)
                    if latin is not None:
                        typeface = latin.get('typeface')
                        if typeface and 'bold' in typeface.lower():
                            font_data['bold'] = True
                    
                    # Color from defRPr
                    solidFill = defRPr.find('a:solidFill', NAMESPACES)
                    if solidFill is not None:
                        srgbClr = solidFill.find('a:srgbClr', NAMESPACES)
                        if srgbClr is not None:
                            val = srgbClr.get('val')
                            if val:
                                font_data['color'] = f"#{val}"
            
            # Get first run for font properties (overrides defRPr)
            r = para.find('.//a:r', NAMESPACES)
            if r is not None:
                rPr = r.find('a:rPr', NAMESPACES)
                if rPr is not None:
                    # Font size (in 1/100 points, convert to CSS pixels)
                    # Formula: (sz/100) * (96/72) = sz/100 * 1.333...
                    sz = rPr.get('sz')
                    if sz:
                        points = int(sz) / 100
                        font_size = round(points * 1.333, 1)
                        # Apply canvas scale factor (same as coordinate scaling)
                        # This ensures fonts are proportional to the scaled canvas
                        # Note: This requires slide_width context which XML extraction doesn't have
                        # For now, fonts from XML extraction won't be pre-scaled
                        # (They'll be manually scaled if needed in Svelte)
                        font_data['font_size'] = font_size
                    
                    # Bold (explicit False overrides defRPr, True overrides, None inherits)
                    b_attr = rPr.get('b')
                    if b_attr == '1':
                        font_data['bold'] = True
                    elif b_attr == '0':
                        font_data['bold'] = False
                    
                    # Italic
                    i_attr = rPr.get('i')
                    if i_attr == '1':
                        font_data['italic'] = True
                    elif i_attr == '0':
                        font_data['italic'] = False

                    # Underline
                    u_attr = rPr.get('u')
                    if u_attr and u_attr != 'none':
                        font_data['underline'] = True
                    elif u_attr == 'none':
                        font_data['underline'] = False
                    
                    # Font name
                    latin = rPr.find('a:latin', NAMESPACES)
                    if latin is not None:
                        typeface = latin.get('typeface')
                        if typeface and 'bold' in typeface.lower():
                            font_data['bold'] = True
                    
                    # Color
                    solidFill = rPr.find('a:solidFill', NAMESPACES)
                    if solidFill is not None:
                        srgbClr = solidFill.find('a:srgbClr', NAMESPACES)
                        if srgbClr is not None:
                            val = srgbClr.get('val')
                            if val:
                                font_data['color'] = f"#{val}"
        
        # Fallback wrap inference from preserved line breaks / overflow heuristic.
        text_content = get_text_from_shape_xml(shape_elem)
        if 'wrap' not in font_data and should_force_font_wrap(text_content, layout, font_data.get('font_size')):
            font_data['wrap'] = True

        return font_data if font_data else None
    except:
        return None

def extract_visual_data_from_xml(shape_elem, z_index, slide_width=None):
    """Extract visual data from shape XML element.
    
    Args:
        shape_elem: XML element (p:sp or p:cxnSp)
        z_index: The z-index for this shape
        slide_width: Source slide width for auto-scaling
    
    Returns:
        Visual data dict
    """
    visual_data = {'z_index': z_index}
    
    # Determine shape type from XML
    if shape_elem.tag.endswith('sp'):
        visual_data['shape_type'] = 'text_box'
    elif shape_elem.tag.endswith('cxnSp'):
        visual_data['shape_type'] = 'connector'
    
    # Extract layout
    layout = extract_layout_from_xml(shape_elem, slide_width)
    if layout:
        visual_data['layout'] = layout
    
    # Extract font
    font = extract_font_from_xml(shape_elem, layout)
    if font:
        visual_data['font'] = font
    
    # TODO: Could also extract fill and line from XML if needed
    # For now, these are less critical for grouped text shapes
    
    return visual_data

def extract_shape_layout(shape, slide_width=None):
    """Extract position, size, and rotation from a shape.
    
    Args:
        shape: The pptx shape object
        slide_width: Source slide width in pixels for auto-scaling (optional)
    
    Returns:
        Layout dict with coordinates scaled to target canvas (960×540)
    """
    try:
        # Extract raw coordinates in source dimensions
        x = emu_to_px(shape.left)
        y = emu_to_px(shape.top)
        width = emu_to_px(shape.width)
        height = emu_to_px(shape.height)
        
        # Auto-scale to target canvas if slide_width provided
        if slide_width is not None:
            x = scale_to_target(x, slide_width)
            y = scale_to_target(y, slide_width)  # Use same scale for y
            width = scale_to_target(width, slide_width)
            height = scale_to_target(height, slide_width)
        
        return {
            'x': x,
            'y': y,
            'width': width,
            'height': height,
            'rotation': shape.rotation if shape.rotation else 0
        }
    except Exception as e:
        return None

def calculate_line_endpoints(layout, slide_width=None):
    """Calculate actual line endpoints from layout with rotation.
    
    PowerPoint stores lines as rectangles with rotation. The line runs from
    top-center to bottom-center of the unrotated rectangle, then the whole
    thing is rotated around the center.
    
    For lines:
    - rotation 0: vertical line (top to bottom)
    - rotation 90 or 270: horizontal line (left to right or right to left)
    - other angles: diagonal line
    
    Args:
        layout: Layout dict with x, y, width, height, rotation (already scaled if slide_width was used)
        slide_width: Not used here - layout is already scaled by extract_shape_layout
    
    Returns dict with 'from' and 'to' points {x, y} or None if not applicable.
    """
    if not layout:
        return None
    
    x = layout.get('x', 0)
    y = layout.get('y', 0)
    w = layout.get('width', 0)
    h = layout.get('height', 0)
    rotation = layout.get('rotation', 0)
    
    # Center of the shape
    cx = x + w / 2
    cy = y + h / 2
    
    # Original endpoints (before rotation) - line from top-center to bottom-center
    # The "height" is the line length in its unrotated state
    p1_x, p1_y = cx, y           # top-center
    p2_x, p2_y = cx, y + h       # bottom-center
    
    # Rotate around center
    rad = math.radians(rotation)
    cos_r = math.cos(rad)
    sin_r = math.sin(rad)
    
    def rotate_point(px, py):
        dx = px - cx
        dy = py - cy
        rx = cx + dx * cos_r - dy * sin_r
        ry = cy + dx * sin_r + dy * cos_r
        return round(rx), round(ry)
    
    from_pt = rotate_point(p1_x, p1_y)
    to_pt = rotate_point(p2_x, p2_y)
    
    return {
        'from': {'x': from_pt[0], 'y': from_pt[1]},
        'to': {'x': to_pt[0], 'y': to_pt[1]}
    }

def extract_fill_style(shape, theme_colors=None):
    """Extract fill color from a shape.
    
    Args:
        shape: The pptx shape object
        theme_colors: Optional dict mapping theme color strings to hex values
    """
    try:
        if not hasattr(shape, 'fill'):
            return None
        fill = shape.fill
        if fill.type is None:
            return None
        
        # Try to get solid fill color
        try:
            fore_color = fill.fore_color
            if fore_color.type is not None:
                # Try RGB first
                try:
                    rgb = fore_color.rgb
                    if rgb:
                        return rgb_to_hex(rgb)
                except:
                    pass
                # Try theme color
                try:
                    theme = fore_color.theme_color
                    brightness = fore_color.brightness
                    theme_str = str(theme)
                    
                    # Resolve theme color to hex
                    if theme_colors and theme_str in theme_colors:
                        hex_color = theme_colors[theme_str]
                        # Return resolved hex (omitting brightness as requested)
                        return hex_color
                    elif theme_colors:
                        # Theme color not in map - warn and use black
                        return resolve_theme_color(theme_str, theme_colors, "fill")
                    else:
                        # Fallback: return theme reference if no theme map provided
                        return {'theme': theme_str, 'brightness': brightness}
                except:
                    pass
        except:
            pass
    except:
        pass
    return None

def extract_line_style(shape, theme_colors=None):
    """Extract line/stroke properties from a shape.
    
    Args:
        shape: The pptx shape object
        theme_colors: Optional dict mapping theme color strings to hex values
    """
    try:
        if not hasattr(shape, 'line'):
            return None
        line = shape.line
        
        line_data = {}
        
        # Get line width
        if line.width:
            line_data['width'] = emu_to_px(line.width)
        
        # Get line color
        try:
            color = line.color
            if color.type is not None:
                try:
                    rgb = color.rgb
                    if rgb:
                        line_data['color'] = rgb_to_hex(rgb)
                except:
                    try:
                        theme = color.theme_color
                        theme_str = str(theme)
                        
                        # Resolve theme color to hex if theme_colors provided
                        if theme_colors and theme_str in theme_colors:
                            line_data['color'] = theme_colors[theme_str]
                        elif theme_colors:
                            # Theme color not found - warn and use black
                            line_data['color'] = resolve_theme_color(theme_str, theme_colors, "line")
                    except:
                        pass
        except:
            pass
        
        # Get dash style from line.dash_style
        try:
            dash_style = line.dash_style
            if dash_style is not None:
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                dash_map = {
                    MSO_LINE_DASH_STYLE.SOLID: None,  # Don't include for solid
                    MSO_LINE_DASH_STYLE.DASH: 'dash',
                    MSO_LINE_DASH_STYLE.DASH_DOT: 'dashDot',
                    MSO_LINE_DASH_STYLE.DASH_DOT_DOT: 'dashDotDot',
                    MSO_LINE_DASH_STYLE.LONG_DASH: 'lgDash',
                    MSO_LINE_DASH_STYLE.LONG_DASH_DOT: 'lgDashDot',
                    MSO_LINE_DASH_STYLE.ROUND_DOT: 'dot',
                    MSO_LINE_DASH_STYLE.SQUARE_DOT: 'sysDash',  # sysDash maps to SQUARE_DOT in python-pptx
                }
                dash_str = dash_map.get(dash_style)
                if dash_str:
                    line_data['dash'] = dash_str
        except:
            pass
        
        if line_data:
            return line_data
    except:
        pass
    return None


def collect_referenced_linked_slides(pptx_path, total_slides, custom_shows):
    """Collect slide numbers referenced by custom shows and slide-jump hyperlinks.

    This augments hidden-slide detection so linked content is captured even when
    a target slide is not marked hidden in PowerPoint.

    Args:
        pptx_path: Path to the PPTX file
        total_slides: Total number of slides in the presentation
        custom_shows: Dict of parsed custom shows

    Returns:
        Set of referenced slide numbers (1-indexed)
    """
    referenced = set()

    # Custom show members are linked content by definition
    for show in custom_shows.values():
        for slide_num in show.get('slide_numbers', []):
            if isinstance(slide_num, int):
                referenced.add(slide_num)

    with zipfile.ZipFile(pptx_path, 'r') as zf:
        for slide_num in range(1, total_slides + 1):
            slide_file = f'ppt/slides/slide{slide_num}.xml'
            try:
                slide_xml = zf.read(slide_file).decode('utf-8')
            except Exception:
                continue

            # Relationship map needed to resolve hlinksldjump rId -> slide number
            rid_to_target_slide = parse_slide_links_from_relationships(pptx_path, slide_num)
            shapes = parse_shapes_from_slide(slide_xml)

            for shape in shapes.values():
                hyperlink = shape.get('hyperlink')
                if not hyperlink:
                    continue

                if hyperlink.get('type') == 'slide':
                    r_id = hyperlink.get('r_id')
                    if r_id and r_id in rid_to_target_slide:
                        referenced.add(rid_to_target_slide[r_id])
                elif hyperlink.get('type') == 'customshow':
                    show_id = hyperlink.get('id')
                    if show_id is not None and show_id in custom_shows:
                        for linked_num in custom_shows[show_id].get('slide_numbers', []):
                            if isinstance(linked_num, int):
                                referenced.add(linked_num)

    return referenced


def extract_arrow_ends_from_xml(slide_xml_content, shape_id):
    """Extract arrow head/tail end types from slide XML.
    
    Returns dict with 'headEnd' and 'tailEnd' if present.
    Values are: 'none', 'triangle', 'stealth', 'diamond', 'oval', 'arrow'
    """
    try:
        root = ET.fromstring(slide_xml_content)
        
        # Find the shape with matching ID (check both sp and cxnSp elements)
        for sp in root.findall('.//p:sp', NAMESPACES):
            cNvPr = sp.find('.//p:cNvPr', NAMESPACES)
            if cNvPr is not None and cNvPr.get('id') == str(shape_id):
                spPr = sp.find('p:spPr', NAMESPACES)
                if spPr is not None:
                    return _extract_line_ends(spPr)
        
        # Also check connector shapes
        for cxnSp in root.findall('.//p:cxnSp', NAMESPACES):
            cNvPr = cxnSp.find('.//p:cNvPr', NAMESPACES)
            if cNvPr is not None and cNvPr.get('id') == str(shape_id):
                spPr = cxnSp.find('p:spPr', NAMESPACES)
                if spPr is not None:
                    return _extract_line_ends(spPr)
    except:
        pass
    return None


def _extract_line_ends(spPr):
    """Extract headEnd and tailEnd from a spPr element."""
    ln = spPr.find('a:ln', NAMESPACES)
    if ln is None:
        return None
    
    result = {}
    
    headEnd = ln.find('a:headEnd', NAMESPACES)
    if headEnd is not None:
        head_type = headEnd.get('type', 'none')
        if head_type and head_type != 'none':
            result['headEnd'] = head_type
    
    tailEnd = ln.find('a:tailEnd', NAMESPACES)
    if tailEnd is not None:
        tail_type = tailEnd.get('type', 'none')
        if tail_type and tail_type != 'none':
            result['tailEnd'] = tail_type
    
    return result if result else None

def extract_font_style(shape, slide_width=None, theme_colors=None, layout=None):
    """Extract font properties from the first text run in a shape.
    
    Args:
        shape: The pptx shape object
        slide_width: Source slide width for auto-scaling font sizes (optional)
        theme_colors: Optional dict mapping theme color strings to hex values
    """
    try:
        if not shape.has_text_frame:
            return None
        
        tf = shape.text_frame
        if not tf.paragraphs:
            return None
        
        font_data = {}
        
        # Get vertical alignment from text frame's vertical_anchor property
        try:
            from pptx.enum.text import MSO_ANCHOR
            v_anchor = tf.vertical_anchor
            if v_anchor == MSO_ANCHOR.TOP:
                font_data['v_align'] = 'top'
            elif v_anchor == MSO_ANCHOR.MIDDLE:
                font_data['v_align'] = 'middle'
            elif v_anchor == MSO_ANCHOR.BOTTOM:
                font_data['v_align'] = 'bottom'
        except:
            pass
        
        # Get text wrapping property from text frame
        try:
            # word_wrap is a boolean property on the text_frame
            # True = text wraps, False = text doesn't wrap
            if hasattr(tf, 'word_wrap') and tf.word_wrap:
                font_data['wrap'] = True
        except:
            pass
        
        # Get horizontal alignment from first paragraph
        para = tf.paragraphs[0]
        if para.alignment:
            from pptx.enum.text import PP_ALIGN
            align_map = {
                PP_ALIGN.LEFT: 'left',
                PP_ALIGN.CENTER: 'center',
                PP_ALIGN.RIGHT: 'right',
            }
            if para.alignment in align_map:
                font_data['align'] = align_map[para.alignment]
        
        # Check paragraph-level default font properties (inherited from theme/master)
        # This captures bold/italic that apply to the entire paragraph
        has_default_bold = False
        has_default_italic = False
        has_default_underline = False
        try:
            # Access the underlying XML to check defRPr (default run properties)
            para_xml = para._element
            pPr = para_xml.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            if pPr is not None:
                defRPr = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
                if defRPr is not None:
                    if defRPr.get('b') == '1':
                        has_default_bold = True
                        font_data['bold'] = True
                    if defRPr.get('i') == '1':
                        has_default_italic = True
                        font_data['italic'] = True
                    if defRPr.get('u') and defRPr.get('u') != 'none':
                        has_default_underline = True
                        font_data['underline'] = True
        except:
            pass
        
        # Check if this is a title placeholder (which is typically bold by theme)
        is_title_placeholder = False
        try:
            if shape.is_placeholder:
                from pptx.enum.shapes import PP_PLACEHOLDER
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    is_title_placeholder = True
        except:
            pass
        
        # Get font properties from first run with text
        found_run = False
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    font = run.font
                    if font.name:
                        # Check if font name includes "Bold" and set bold property
                        if 'bold' in font.name.lower():
                            font_data['bold'] = True
                    if font.size:
                        # font.size.pt is in PowerPoint points (1/72 inch)
                        # CSS pixels are at 96 DPI, so conversion is: points × (96/72) = points × 1.333...
                        # This is the standard DPI conversion formula and is correct.
                        font_data['font_size'] = round(font.size.pt * (96 / 72), 1)
                        # Apply canvas scale factor (same as coordinate scaling)
                        # This ensures fonts are proportional to the scaled canvas
                        if slide_width is not None:
                            scale_factor = TARGET_CANVAS_WIDTH / slide_width
                            font_data['font_size'] = round(font_data['font_size'] * scale_factor, 1)
                    # Bold - explicitly check for True (overrides default), False (overrides default), or None (uses default/theme)
                    if font.bold is True:
                        font_data['bold'] = True
                    elif font.bold is False:
                        # Explicitly not bold - overrides any default
                        if 'bold' in font_data:
                            del font_data['bold']
                    elif font.bold is None and is_title_placeholder and 'bold' not in font_data:
                        # Title placeholders are typically bold by theme
                        font_data['bold'] = True
                    # If font.bold is None and we already have default bold, keep it
                    
                    # Italic - same logic
                    if font.italic is True:
                        font_data['italic'] = True
                    elif font.italic is False:
                        # Explicitly not italic - overrides any default
                        if 'italic' in font_data:
                            del font_data['italic']

                    # Underline - same logic
                    if font.underline is True:
                        font_data['underline'] = True
                    elif font.underline is False:
                        if 'underline' in font_data:
                            del font_data['underline']
                    elif font.underline is None and has_default_underline and 'underline' not in font_data:
                        font_data['underline'] = True
                    
                    # Get font color
                    try:
                        fc = font.color
                        if fc.type is not None:
                            try:
                                rgb = fc.rgb
                                if rgb:
                                    font_data['color'] = rgb_to_hex(rgb)
                            except:
                                try:
                                    theme = fc.theme_color
                                    theme_str = str(theme)
                                    
                                    # Resolve theme color to hex
                                    if theme_colors and theme_str in theme_colors:
                                        font_data['color'] = theme_colors[theme_str]
                                    elif theme_colors:
                                        font_data['color'] = resolve_theme_color(theme_str, theme_colors, "font")
                                    else:
                                        # No theme_colors provided - include theme_color reference
                                        font_data['theme_color'] = theme_str
                                except:
                                    pass
                    except:
                        pass
                    
                    found_run = True
                    break
            if found_run:
                break

        # Fallback wrap inference:
        # - Manual line breaks in text (paragraph returns / soft breaks)
        # - Long text overflow relative to box width and font size
        if 'wrap' not in font_data:
            shape_text = ''
            try:
                shape_text = shape.text or ''
            except:
                shape_text = ''
            if should_force_font_wrap(shape_text, layout, font_data.get('font_size')):
                font_data['wrap'] = True
        
        return font_data if font_data else None
    except:
        return None

def get_shape_type_name(shape):
    """Get a simplified shape type name."""
    try:
        shape_type = shape.shape_type
        if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Try to get the auto shape type
            try:
                auto_type = shape.auto_shape_type
                auto_str = str(auto_type).replace('MSO_AUTO_SHAPE_TYPE.', '').lower()
                if 'arrow' in auto_str:
                    return 'arrow'
                return auto_str
            except:
                return 'auto_shape'
        elif shape_type == MSO_SHAPE_TYPE.LINE:
            return 'connector'
        elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return 'text_box'
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            return 'picture'
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            return 'group'
        elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return 'freeform'
        else:
            return str(shape_type).replace('MSO_SHAPE_TYPE.', '').lower()
    except:
        return 'unknown'

def extract_connector_path(shape):
    """Extract start and end points for connector shapes."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            # For lines, calculate start and end from position and size
            x = emu_to_px(shape.left)
            y = emu_to_px(shape.top)
            w = emu_to_px(shape.width)
            h = emu_to_px(shape.height)
            
            # Determine direction based on the shape's flip properties
            return {
                'start': {'x': x, 'y': y},
                'end': {'x': x + w, 'y': y + h}
            }
    except:
        pass
    return None

def extract_arc_path_from_xml(slide_xml_content, shape_id, shape_layout, slide_width=None):
    """Extract arc path data for a freeform shape from slide XML.
    
    Args:
        slide_xml_content: The slide XML content
        shape_id: The shape ID to extract
        shape_layout: The shape layout (already scaled if slide_width was used)
        slide_width: Not used here - shape_layout is already scaled
    
    Returns arc parameters:
    - from: start point in canvas coordinates (scaled)
    - to: end point in canvas coordinates (scaled)
    - curve: vertical offset for quadratic bezier (negative = up, positive = down)
    - flip: whether the arc is horizontally flipped
    """
    try:
        root = ET.fromstring(slide_xml_content)
        
        # Find the shape with matching ID
        for sp in root.findall('.//p:sp', NAMESPACES):
            cNvPr = sp.find('.//p:cNvPr', NAMESPACES)
            if cNvPr is not None and cNvPr.get('id') == str(shape_id):
                # Check if it's a freeform with custom geometry
                spPr = sp.find('p:spPr', NAMESPACES)
                if spPr is None:
                    continue
                
                custGeom = spPr.find('a:custGeom', NAMESPACES)
                if custGeom is None:
                    continue
                
                # Get transform info
                xfrm = spPr.find('a:xfrm', NAMESPACES)
                flipH = xfrm.get('flipH') == '1' if xfrm is not None else False
                flipV = xfrm.get('flipV') == '1' if xfrm is not None else False
                
                # Find the first path (the stroke path, not fill path)
                pathLst = custGeom.find('a:pathLst', NAMESPACES)
                if pathLst is None:
                    continue
                
                # Get the path that has fill="none" (stroke path)
                stroke_path = None
                for path in pathLst.findall('a:path', NAMESPACES):
                    if path.get('fill') == 'none':
                        stroke_path = path
                        break
                
                if stroke_path is None:
                    # Fall back to first path
                    stroke_path = pathLst.find('a:path', NAMESPACES)
                
                if stroke_path is None:
                    continue
                
                # Get path dimensions for coordinate scaling
                path_w = int(stroke_path.get('w', '21600'))
                path_h = int(stroke_path.get('h', '21600'))
                
                # Extract all points from the path
                points = []
                
                # moveTo is the start point
                moveTo = stroke_path.find('a:moveTo', NAMESPACES)
                if moveTo is not None:
                    pt = moveTo.find('a:pt', NAMESPACES)
                    if pt is not None:
                        points.append({
                            'type': 'move',
                            'x': int(pt.get('x', '0')),
                            'y': int(pt.get('y', '0'))
                        })
                
                # cubicBezTo contains control and end points
                for bezier in stroke_path.findall('a:cubicBezTo', NAMESPACES):
                    pts = bezier.findall('a:pt', NAMESPACES)
                    if len(pts) >= 3:
                        # First two are control points, third is end point
                        points.append({
                            'type': 'cubic',
                            'cp1_x': int(pts[0].get('x', '0')),
                            'cp1_y': int(pts[0].get('y', '0')),
                            'cp2_x': int(pts[1].get('x', '0')),
                            'cp2_y': int(pts[1].get('y', '0')),
                            'x': int(pts[2].get('x', '0')),
                            'y': int(pts[2].get('y', '0'))
                        })
                
                if len(points) < 2:
                    continue
                
                # Get shape layout for coordinate conversion
                layout_x = shape_layout.get('x', 0)
                layout_y = shape_layout.get('y', 0)
                layout_w = shape_layout.get('width', 100)
                layout_h = shape_layout.get('height', 50)
                
                # Scale path coordinates to canvas coordinates
                def scale_x(px):
                    scaled = (px / path_w) * layout_w
                    if flipH:
                        scaled = layout_w - scaled
                    return round(layout_x + scaled)
                
                def scale_y(py):
                    scaled = (py / path_h) * layout_h
                    if flipV:
                        scaled = layout_h - scaled
                    return round(layout_y + scaled)
                
                # Get start and end points
                start_point = points[0]
                from_x = scale_x(start_point.get('x', 0))
                from_y = scale_y(start_point.get('y', 0))
                
                # Find the last endpoint
                end_point = None
                for p in reversed(points):
                    if p['type'] == 'cubic':
                        end_point = p
                        break
                
                if end_point is None:
                    continue
                
                to_x = scale_x(end_point.get('x', 0))
                to_y = scale_y(end_point.get('y', 0))
                
                # Calculate curve amount based on control points
                # For a typical arc, we want the midpoint's vertical offset
                # Estimate from the layout height and whether it curves up or down
                mid_y = (from_y + to_y) / 2
                
                # Check if the control points curve up or down
                # A typical arc has control points either above or below the endpoints
                curve_direction = -1  # default: curve up
                if len(points) > 1 and points[1]['type'] == 'cubic':
                    cp1_y = scale_y(points[1].get('cp1_y', 0))
                    # If control point is below the endpoints, curve is down
                    if cp1_y > mid_y:
                        curve_direction = 1
                
                # Curve amount is approximately the height of the arc
                curve_amount = layout_h * curve_direction * 0.8
                
                return {
                    'from': {'x': from_x, 'y': from_y},
                    'to': {'x': to_x, 'y': to_y},
                    'curve': round(curve_amount)
                }
                
    except Exception as e:
        # Silently fail for shapes without arc data
        pass
    
    return None


def extract_images_from_pptx(pptx_path, output_folder):
    """Extract all images from a PPTX file to a folder.
    
    Args:
        pptx_path: Path to the PPTX file
        output_folder: Path to the output folder for images
    
    Returns:
        Dict mapping image paths (e.g., 'media/image1.jpg') to saved filenames
    """
    image_map = {}
    
    # Create output folder if it doesn't exist
    Path(output_folder).mkdir(parents=True, exist_ok=True)
    
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        # Find all image files in the media folder
        for file_info in zf.infolist():
            if file_info.filename.startswith('ppt/media/'):
                # Extract filename
                image_filename = os.path.basename(file_info.filename)
                output_path = os.path.join(output_folder, image_filename)
                
                # Extract image
                with zf.open(file_info.filename) as src:
                    with open(output_path, 'wb') as dst:
                        dst.write(src.read())
                
                # Map the media path to filename (e.g., '../media/image1.jpg' -> 'image1.jpg')
                media_path = file_info.filename.replace('ppt/', '')
                image_map[media_path] = image_filename
                print(f"Extracted image: {image_filename}")
    
    return image_map


def parse_slide_relationships(pptx_path, slide_num):
    """Parse relationship file for a slide to map rId to image paths.
    
    Args:
        pptx_path: Path to the PPTX file
        slide_num: Slide number (1-indexed)
    
    Returns:
        Dict mapping rId to image filename (e.g., 'rId3' -> 'image1.jpg')
    """
    rid_to_image = {}
    rels_file = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
    
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            rels_xml = zf.read(rels_file).decode('utf-8')
            root = ET.fromstring(rels_xml)
            
            # Namespace for relationships
            rels_ns = '{http://schemas.openxmlformats.org/package/2006/relationships}'
            
            for rel in root.findall(f'{rels_ns}Relationship'):
                rel_type = rel.get('Type', '')
                if 'image' in rel_type:
                    rid = rel.get('Id')
                    target = rel.get('Target', '')
                    # Target is like '../media/image1.jpg'
                    image_filename = os.path.basename(target)
                    if rid:
                        rid_to_image[rid] = image_filename
        except KeyError:
            pass  # No relationships file for this slide
    
    return rid_to_image


def parse_slide_links_from_relationships(pptx_path, slide_num):
    """Parse relationship file for a slide to map rId to target slide numbers.
    
    Args:
        pptx_path: Path to the PPTX file
        slide_num: Slide number (1-indexed)
    
    Returns:
        Dict mapping rId to target slide number (e.g., 'rId3' -> 4)
    """
    rid_to_slide = {}
    rels_file = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
    
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            rels_xml = zf.read(rels_file).decode('utf-8')
            root = ET.fromstring(rels_xml)
            
            # Namespace for relationships
            rels_ns = '{http://schemas.openxmlformats.org/package/2006/relationships}'
            
            for rel in root.findall(f'{rels_ns}Relationship'):
                rel_type = rel.get('Type', '')
                # Check for slide relationship type
                if rel_type.endswith('/slide'):
                    rid = rel.get('Id')
                    target = rel.get('Target', '')
                    # Target is like 'slide4.xml'
                    match = re.search(r'slide(\d+)\.xml', target)
                    if match and rid:
                        rid_to_slide[rid] = int(match.group(1))
        except KeyError:
            pass  # No relationships file for this slide
    
    return rid_to_slide


def parse_pictures_from_slide(slide_xml_content, rid_to_image, slide_width=None, shape_z_indices=None):
    """Parse all picture elements from slide XML.
    
    Args:
        slide_xml_content: The slide XML content
        rid_to_image: Dict mapping rId to image filenames
        slide_width: Source slide width for auto-scaling
        shape_z_indices: Optional dict mapping shape ID to z_index
    
    Returns:
        Dict mapping shape ID to picture data
    """
    root = ET.fromstring(slide_xml_content)
    pictures = {}
    
    # Find all p:pic elements
    for pic in root.findall('.//p:pic', NAMESPACES):
        # Get shape ID and name from nvPicPr
        nvPicPr = pic.find('p:nvPicPr', NAMESPACES)
        if nvPicPr is None:
            continue
        
        cNvPr = nvPicPr.find('p:cNvPr', NAMESPACES)
        if cNvPr is None:
            continue
        
        shape_id = cNvPr.get('id')
        shape_name = cNvPr.get('name', '')
        description = cNvPr.get('descr', '')  # Original filename often stored here
        
        if not shape_id:
            continue
        
        # Get image reference from blipFill
        blipFill = pic.find('p:blipFill', NAMESPACES)
        if blipFill is None:
            continue
        
        blip = blipFill.find('a:blip', NAMESPACES)
        if blip is None:
            continue
        
        # Get the relationship ID (r:embed attribute)
        embed_rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
        if not embed_rid:
            continue
        
        # Look up the image filename
        image_filename = rid_to_image.get(embed_rid)
        if not image_filename:
            continue
        
        # Get layout from spPr
        spPr = pic.find('p:spPr', NAMESPACES)
        layout = None
        line_style = None
        
        if spPr is not None:
            xfrm = spPr.find('a:xfrm', NAMESPACES)
            if xfrm is not None:
                off = xfrm.find('a:off', NAMESPACES)
                ext = xfrm.find('a:ext', NAMESPACES)
                
                if off is not None and ext is not None:
                    x = emu_to_px(int(off.get('x', 0)))
                    y = emu_to_px(int(off.get('y', 0)))
                    width = emu_to_px(int(ext.get('cx', 0)))
                    height = emu_to_px(int(ext.get('cy', 0)))
                    rotation = int(xfrm.get('rot', 0)) / 60000
                    
                    # Scale to target canvas
                    if slide_width is not None:
                        x = scale_to_target(x, slide_width)
                        y = scale_to_target(y, slide_width)
                        width = scale_to_target(width, slide_width)
                        height = scale_to_target(height, slide_width)
                    
                    layout = {
                        'x': x,
                        'y': y,
                        'width': width,
                        'height': height,
                        'rotation': rotation
                    }
            
            # Check for line/border style
            ln = spPr.find('a:ln', NAMESPACES)
            if ln is not None:
                line_style = {}
                width_attr = ln.get('w')
                if width_attr:
                    line_style['width'] = emu_to_px(int(width_attr))
                
                solidFill = ln.find('a:solidFill', NAMESPACES)
                if solidFill is not None:
                    srgbClr = solidFill.find('a:srgbClr', NAMESPACES)
                    if srgbClr is not None:
                        line_style['color'] = f"#{srgbClr.get('val', '000000')}"
                    else:
                        schemeClr = solidFill.find('a:schemeClr', NAMESPACES)
                        if schemeClr is not None:
                            line_style['theme_color'] = schemeClr.get('val')
        
        # Get z_index from shape_z_indices if provided
        z_index = shape_z_indices.get(shape_id) if shape_z_indices else None
        
        pictures[shape_id] = {
            'id': shape_id,
            'name': shape_name,
            'image': image_filename,
            'description': description,
            'z_index': z_index,
            'layout': layout,
            'line': line_style if line_style else None
        }
    
    return pictures


def extract_shape_visual_data(shape, z_index, slide_xml_content=None, shape_id=None, slide_width=None, theme_colors=None):
    """Extract all visual data for a shape.
    
    Args:
        shape: The pptx shape object
        z_index: The z-index of the shape
        slide_xml_content: Optional slide XML for extracting arrow/arc data
        shape_id: Optional shape ID for XML lookups
        slide_width: Source slide width for auto-scaling coordinates
        theme_colors: Optional dict mapping theme color strings to hex values
    """
    visual_data = {
        'z_index': z_index,
        'shape_type': get_shape_type_name(shape)
    }
    
    # Layout (position, size, rotation) - auto-scaled to target canvas
    layout = extract_shape_layout(shape, slide_width)
    if layout:
        visual_data['layout'] = layout
    
    # Fill color
    fill = extract_fill_style(shape, theme_colors)
    if fill:
        visual_data['fill'] = fill
    
    # Line/stroke
    line = extract_line_style(shape, theme_colors)
    if line:
        visual_data['line'] = line
    
    # Font properties
    font = extract_font_style(shape, slide_width, theme_colors, layout)
    if font:
        visual_data['font'] = font
    
    # Connector path (for lines/arrows)
    # Check shape_type OR shape.name containing "Line" (some lines are auto_shape type)
    is_line_shape = visual_data['shape_type'] in ('connector', 'line')
    try:
        if hasattr(shape, 'name') and shape.name and 'line' in shape.name.lower():
            is_line_shape = True
    except:
        pass
    
    if is_line_shape:
        path = extract_connector_path(shape)
        if path:
            visual_data['path'] = path
        
        # Calculate actual line endpoints from layout + rotation
        # Note: layout is already scaled, so no need to pass slide_width
        if layout:
            line_endpoints = calculate_line_endpoints(layout)
            if line_endpoints:
                visual_data['line_endpoints'] = line_endpoints
    
    # Arrow head/tail ends from XML (for lines that are actually arrows)
    if slide_xml_content and shape_id:
        arrow_ends = extract_arrow_ends_from_xml(slide_xml_content, shape_id)
        if arrow_ends:
            visual_data['arrow_ends'] = arrow_ends
    
    # Arc path (for freeform arcs)
    # Note: layout is already scaled, so no need to pass slide_width
    if visual_data['shape_type'] == 'freeform' and slide_xml_content and shape_id and layout:
        arc_path = extract_arc_path_from_xml(slide_xml_content, shape_id, layout)
        if arc_path:
            visual_data['arc_path'] = arc_path
    
    return visual_data

def get_text_from_shape_xml(shape_elem):
    """Extract all text from a shape XML element, preserving line and paragraph breaks."""
    tx_body = shape_elem.find('.//p:txBody', NAMESPACES)
    if tx_body is None:
        return ''

    paragraph_texts = []
    for para in tx_body.findall('a:p', NAMESPACES):
        parts = []
        for node in list(para):
            # Regular text run
            if node.tag.endswith('}r'):
                text_node = node.find('a:t', NAMESPACES)
                if text_node is not None and text_node.text:
                    parts.append(text_node.text)
            # Explicit line break within paragraph
            elif node.tag.endswith('}br'):
                parts.append('\n')
            # Text field run
            elif node.tag.endswith('}fld'):
                text_node = node.find('a:t', NAMESPACES)
                if text_node is not None and text_node.text:
                    parts.append(text_node.text)

        paragraph_texts.append(''.join(parts))

    return normalize_unicode_text('\n'.join(paragraph_texts).strip())


def _extract_run_format_from_rpr(rPr, base_format=None):
    """Extract run formatting from rPr, optionally applying overrides to base_format."""
    fmt = dict(base_format) if base_format else {}
    if rPr is None:
        return fmt

    b_attr = rPr.get('b')
    if b_attr == '1':
        fmt['bold'] = True
    elif b_attr == '0':
        fmt['bold'] = False

    i_attr = rPr.get('i')
    if i_attr == '1':
        fmt['italic'] = True
    elif i_attr == '0':
        fmt['italic'] = False

    u_attr = rPr.get('u')
    if u_attr:
        fmt['underline'] = (u_attr != 'none')

    strike_attr = rPr.get('strike')
    if strike_attr:
        fmt['strike'] = (strike_attr != 'noStrike')

    baseline_attr = rPr.get('baseline')
    if baseline_attr is not None:
        try:
            baseline = int(baseline_attr)
            if baseline > 0:
                fmt['superscript'] = True
                fmt['subscript'] = False
            elif baseline < 0:
                fmt['subscript'] = True
                fmt['superscript'] = False
            else:
                fmt['superscript'] = False
                fmt['subscript'] = False
        except ValueError:
            pass

    return fmt


def _format_run_as_html(text, fmt, scripture_mode=False):
    """Format a run's text as Fragment-friendly minimal HTML and indicate formatting."""
    escaped = html.escape(text, quote=False)
    if not escaped:
        return escaped, False

    has_formatting = False
    content = escaped
    has_scripture_marker = False

    if fmt.get('subscript'):
        has_formatting = True
        has_scripture_marker = True
        content = f"<sub>{content}</sub>"
    elif fmt.get('superscript'):
        has_formatting = True
        has_scripture_marker = True
        content = f"<sup>{content}</sup>"

    if scripture_mode and has_scripture_marker:
        return content, has_formatting

    if fmt.get('strike'):
        has_formatting = True
        content = f"<s>{content}</s>"
    if fmt.get('underline'):
        has_formatting = True
        content = f"<u>{content}</u>"
    if fmt.get('italic'):
        has_formatting = True
        content = f"<em>{content}</em>"
    if fmt.get('bold'):
        has_formatting = True
        content = f"<strong>{content}</strong>"

    return content, has_formatting


def _normalize_fragment_markup(text):
    """Normalize generated inline markup by collapsing adjacent identical tags.

    PowerPoint frequently splits visually continuous text into multiple runs,
    which can produce sequences like </strong><strong> around punctuation.
    """
    if not text:
        return text

    normalized = text
    # Only collapse simple wrapper tags used by extractor.
    for tag in ('strong', 'em', 'u', 's'):
        normalized = normalized.replace(f'</{tag}><{tag}>', '')

    # Drop empty wrappers that can occur after run collapsing.
    normalized = re.sub(r'<(strong|em|u|s)\s*>\s*</\1>', '', normalized)

    # Canonicalize all break tag variants.
    normalized = re.sub(r'<br\s*/?>', '<br/>', normalized)

    # Preserve intentional blank lines (<br/><br/>) but collapse accidental longer runs.
    normalized = re.sub(r'(?:<br\s*/?>\s*){3,}', '<br/><br/>', normalized)

    return normalized


def get_text_with_fragment_markup_from_shape_xml(shape_elem):
    """Extract plain text plus optional minimal HTML markup from shape XML.

    Returns:
        tuple[str, str, bool, bool]: (plain_text, text_with_markup, has_markup, is_scripture)
    """
    tx_body = shape_elem.find('.//p:txBody', NAMESPACES)
    if tx_body is None:
        return '', '', False, False

    paragraph_texts = []
    paragraph_html = []
    has_rich_formatting = False
    scripture_mode = False

    # Heuristic A: leading "{Book} {Chapter}:{VerseOrRange}" + body text.
    # This catches scripture blocks that do not include superscript verse numbers.
    raw_text = get_text_from_shape_xml(shape_elem)
    if _looks_like_scripture_reference_with_text(raw_text):
        scripture_mode = True

    # Scripture detection heuristic (current):
    # - A) leading reference + body text pattern (see above)
    # - B) any superscript run marks this block as scripture
    #
    # Where this comes from:
    # - _extract_run_format_from_rpr reads a:rPr@baseline
    # - baseline > 0 => run_fmt['superscript'] = True
    # - if any run/fld in the text body is superscript, we set scripture_mode = True
    # These are intentional best-effort signals and may evolve.
    #
    # Important: this heuristic may evolve, but downstream behavior should remain explicit:
    # parse_shapes_from_slide stores `is_scripture`, and entry assembly applies
    # _sanitize_scripture_block_font(entry, is_scripture=...).
    #
    # Keep this detection logic and downstream sanitation contract in sync.
    for para in tx_body.findall('a:p', NAMESPACES):
        paragraph_default_format = {}
        pPr = para.find('a:pPr', NAMESPACES)
        if pPr is not None:
            defRPr = pPr.find('a:defRPr', NAMESPACES)
            paragraph_default_format = _extract_run_format_from_rpr(defRPr)

        for node in list(para):
            if node.tag.endswith('}r') or node.tag.endswith('}fld'):
                rPr = node.find('a:rPr', NAMESPACES)
                run_fmt = _extract_run_format_from_rpr(rPr, paragraph_default_format)
                if run_fmt.get('superscript'):
                    scripture_mode = True
                    break
        if scripture_mode:
            break

    for para in tx_body.findall('a:p', NAMESPACES):
        parts = []
        html_parts = []

        paragraph_default_format = {}
        pPr = para.find('a:pPr', NAMESPACES)
        if pPr is not None:
            defRPr = pPr.find('a:defRPr', NAMESPACES)
            paragraph_default_format = _extract_run_format_from_rpr(defRPr)

        for node in list(para):
            # Regular text run
            if node.tag.endswith('}r'):
                text_node = node.find('a:t', NAMESPACES)
                if text_node is not None and text_node.text:
                    run_text = normalize_unicode_text(text_node.text)
                    parts.append(run_text)
                    rPr = node.find('a:rPr', NAMESPACES)
                    run_fmt = _extract_run_format_from_rpr(rPr, paragraph_default_format)
                    run_html, run_has_formatting = _format_run_as_html(run_text, run_fmt, scripture_mode)
                    html_parts.append(run_html)
                    if run_has_formatting:
                        has_rich_formatting = True
            # Explicit line break within paragraph
            elif node.tag.endswith('}br'):
                parts.append('\n')
                html_parts.append('<br/>')
            # Text field run
            elif node.tag.endswith('}fld'):
                text_node = node.find('a:t', NAMESPACES)
                if text_node is not None and text_node.text:
                    run_text = normalize_unicode_text(text_node.text)
                    parts.append(run_text)
                    rPr = node.find('a:rPr', NAMESPACES)
                    run_fmt = _extract_run_format_from_rpr(rPr, paragraph_default_format)
                    run_html, run_has_formatting = _format_run_as_html(run_text, run_fmt, scripture_mode)
                    html_parts.append(run_html)
                    if run_has_formatting:
                        has_rich_formatting = True

        paragraph_texts.append(''.join(parts))
        paragraph_html.append(''.join(html_parts))

    plain_text = '\n'.join(paragraph_texts).strip()
    text_with_markup = '<br/>'.join(paragraph_html).strip()
    text_with_markup = _normalize_fragment_markup(text_with_markup)

    if has_rich_formatting:
        return plain_text, text_with_markup, True, scripture_mode
    return plain_text, plain_text, False, scripture_mode

def get_hyperlink_from_shape_xml(shape_elem):
    """Extract hyperlink action from shape XML element.
    
    Returns dict with hyperlink info:
    - For custom shows: {'type': 'customshow', 'id': <show_id>}
    - For slide jumps: {'type': 'slide', 'r_id': <relationship_id>}
    - For other actions: {'type': 'action', 'action': <action_string>}
    """
    # Check for hlinkClick on the shape itself (cNvPr)
    cNvPr = shape_elem.find('.//p:cNvPr', NAMESPACES)
    if cNvPr is not None:
        hlinkClick = cNvPr.find('a:hlinkClick', NAMESPACES)
        if hlinkClick is not None:
            action = hlinkClick.get('action', '')
            r_id = hlinkClick.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
            
            if 'customshow' in action.lower():
                # Extract custom show ID
                import re
                match = re.search(r'id=(\d+)', action)
                if match:
                    return {'type': 'customshow', 'id': int(match.group(1))}
            elif 'hlinksldjump' in action.lower() and r_id:
                # Slide jump - return the relationship ID to be resolved later
                return {'type': 'slide', 'r_id': r_id}
            elif action:
                return {'type': 'action', 'action': action}
    
    # Check for hyperlinks in text runs
    for hlinkClick in shape_elem.findall('.//a:hlinkClick', NAMESPACES):
        action = hlinkClick.get('action', '')
        r_id = hlinkClick.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
        
        if 'customshow' in action.lower():
            import re
            match = re.search(r'id=(\d+)', action)
            if match:
                return {'type': 'customshow', 'id': int(match.group(1))}
        elif 'hlinksldjump' in action.lower() and r_id:
            # Slide jump - return the relationship ID to be resolved later
            return {'type': 'slide', 'r_id': r_id}
        elif action:
            return {'type': 'action', 'action': action}
    
    return None

def get_group_child_ids(slide_xml_content, group_id):
    """Get all child shape IDs from a group shape.
    
    Args:
        slide_xml_content: The slide XML content
        group_id: The group shape ID
    
    Returns:
        List of child shape IDs, or None if not a group
    """
    try:
        root = ET.fromstring(slide_xml_content)
        
        # Find the group shape with matching ID
        for grpSp in root.findall('.//p:grpSp', NAMESPACES):
            cNvPr = grpSp.find('.//p:cNvPr', NAMESPACES)
            if cNvPr is not None and cNvPr.get('id') == str(group_id):
                # Found the group - get all child shape IDs
                child_ids = []
                # Get direct child shapes (sp elements)
                for sp in grpSp.findall('./p:sp', NAMESPACES):
                    child_cNvPr = sp.find('.//p:cNvPr', NAMESPACES)
                    if child_cNvPr is not None:
                        child_ids.append(child_cNvPr.get('id'))
                # Also get child connectors (cxnSp elements)
                for cxnSp in grpSp.findall('./p:cxnSp', NAMESPACES):
                    child_cNvPr = cxnSp.find('.//p:cNvPr', NAMESPACES)
                    if child_cNvPr is not None:
                        child_ids.append(child_cNvPr.get('id'))
                return child_ids if child_ids else None
    except:
        pass
    return None

def parse_animation_sequence(slide_xml_content):
    """Parse animation sequence from slide XML and return ordered list of animation entries.
    
    Each entry contains:
    - shape_id: The shape's ID (or list of IDs if it's a group)
    - is_group: Boolean indicating if this is a group animation
    - timing: 'click' (On Click), 'with' (With Previous), or 'after' (After Previous)
    - delay: Delay in milliseconds (for 'after' timing)
    
    The animation structure in PowerPoint XML is:
    - p:timing > p:tnLst > p:par > p:cTn[nodeType="tmRoot"]
    - Inside that: p:seq > p:cTn[nodeType="mainSeq"] > p:childTnLst
    - Each click group: p:par > p:cTn[nodeType="clickPar"]
    - Inside clickPar: p:par > p:cTn[nodeType="withGroup"] or p:cTn[nodeType="afterGroup"]
    - Individual animations: p:cTn[nodeType="clickEffect"|"withEffect"|"afterEffect"]
    - Shape target: p:spTgt spid="..."
    """
    root = ET.fromstring(slide_xml_content)
    animation_entries = []
    seen_shapes = set()
    
    # Find the main sequence
    main_seq = root.find('.//p:cTn[@nodeType="mainSeq"]', NAMESPACES)
    if main_seq is None:
        # Fallback: return empty list if no animations
        return []
    
    # Get the childTnLst which contains click groups
    child_list = main_seq.find('p:childTnLst', NAMESPACES)
    if child_list is None:
        return []
    
    # Iterate through click groups (each p:par with clickPar)
    for click_group in child_list.findall('p:par', NAMESPACES):
        # Process all animations within this click group
        _process_animation_group(click_group, animation_entries, seen_shapes, 0)
    
    return animation_entries


def _process_animation_group(group_elem, entries, seen_shapes, parent_delay):
    """Recursively process animation groups to extract shape timing info.
    
    Args:
        group_elem: The p:par element to process
        entries: List to append animation entries to
        seen_shapes: Set of already-seen shape IDs (to avoid duplicates)
        parent_delay: Accumulated delay from parent afterGroup elements (in ms)
    """
    cTn = group_elem.find('p:cTn', NAMESPACES)
    if cTn is None:
        return
    
    node_type = cTn.get('nodeType', '')
    
    # Get delay from this element's stCondLst if present
    local_delay = 0
    stCondLst = cTn.find('p:stCondLst', NAMESPACES)
    if stCondLst is not None:
        cond = stCondLst.find('p:cond', NAMESPACES)
        if cond is not None:
            delay_str = cond.get('delay', '')
            if delay_str and delay_str != 'indefinite':
                try:
                    local_delay = int(delay_str)
                except ValueError:
                    pass
    
    # Determine timing type from nodeType
    timing = None
    if node_type == 'clickEffect':
        timing = 'click'
    elif node_type == 'withEffect':
        timing = 'with'
    elif node_type == 'afterEffect':
        timing = 'after'
    
    # If this is an animation effect, find the target shape
    if timing is not None:
        # Look for spTgt inside this element
        for spTgt in cTn.findall('.//p:spTgt', NAMESPACES):
            spid = spTgt.get('spid')
            if spid and spid not in seen_shapes:
                seen_shapes.add(spid)
                entry = {
                    'shape_id': spid,
                    'timing': timing
                }
                # Add delay for afterEffect (parent_delay from afterGroup + any local delay)
                total_delay = parent_delay + local_delay
                if timing == 'after' or total_delay > 0:
                    entry['delay'] = total_delay
                entries.append(entry)
                break  # One shape per animation effect
    
    # For afterGroup, accumulate the delay for child animations
    accumulated_delay = parent_delay
    if node_type == 'afterGroup':
        accumulated_delay = parent_delay + local_delay
    
    # Recursively process child elements
    child_list = cTn.find('p:childTnLst', NAMESPACES)
    if child_list is not None:
        for child_par in child_list.findall('p:par', NAMESPACES):
            _process_animation_group(child_par, entries, seen_shapes, accumulated_delay)

def parse_shapes_from_slide(slide_xml_content):
    """Parse all shapes from slide XML and return dict keyed by shape ID."""
    root = ET.fromstring(slide_xml_content)
    shapes = {}
    
    # Find all sp (shape) elements
    for sp in root.findall('.//p:sp', NAMESPACES):
        nvSpPr = sp.find('p:nvSpPr', NAMESPACES)
        if nvSpPr is not None:
            cNvPr = nvSpPr.find('p:cNvPr', NAMESPACES)
            if cNvPr is not None:
                shape_id = cNvPr.get('id')
                shape_name = cNvPr.get('name', '')
                text_plain, text_markup, has_markup, is_scripture = get_text_with_fragment_markup_from_shape_xml(sp)
                hyperlink = get_hyperlink_from_shape_xml(sp)
                
                if shape_id:
                    shape_data = {
                        'id': shape_id,
                        'name': shape_name,
                        'text': text_markup if has_markup else text_plain,
                        'hyperlink': hyperlink,
                        'is_scripture': is_scripture
                    }
                    shapes[shape_id] = shape_data
    
    # Also find connector/line shapes (cxnSp elements)
    for cxn in root.findall('.//p:cxnSp', NAMESPACES):
        nvCxnSpPr = cxn.find('p:nvCxnSpPr', NAMESPACES)
        if nvCxnSpPr is not None:
            cNvPr = nvCxnSpPr.find('p:cNvPr', NAMESPACES)
            if cNvPr is not None:
                shape_id = cNvPr.get('id')
                shape_name = cNvPr.get('name', '')
                
                if shape_id:
                    shapes[shape_id] = {
                        'id': shape_id,
                        'name': shape_name,
                        'text': '',
                        'hyperlink': None,
                        'is_scripture': False,
                        'is_connector': True
                    }
    
    return shapes

def get_hidden_slides(pptx_path, total_slides):
    """Get set of hidden slide numbers by checking each slide's XML file.
    
    Hidden slides have show="0" attribute on the root p:sld element in their
    individual slide XML files (ppt/slides/slideN.xml).
    These are the slides that should go into linked_slides, not the main slides array.
    
    Args:
        pptx_path: Path to the PPTX file
        total_slides: Total number of slides in the presentation
    
    Returns:
        Set of hidden slide numbers (1-indexed)
    """
    hidden_slides = set()
    
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        for slide_num in range(1, total_slides + 1):
            try:
                slide_file = f'ppt/slides/slide{slide_num}.xml'
                slide_xml = zf.read(slide_file).decode('utf-8')
                root = ET.fromstring(slide_xml)
                
                # Check show attribute on root p:sld element
                # show="0" means hidden, show="1" or absent means visible
                show_attr = root.get('show', '1')
                if show_attr == '0':
                    hidden_slides.add(slide_num)
            except Exception as e:
                print(f"Warning: Could not check slide {slide_num} hidden status: {e}")
    
    return hidden_slides


def parse_custom_shows(pptx_path):
    """Parse custom shows from presentation.xml - returns metadata only (not slide content).
    
    Custom shows are named collections of slides that can be linked from the main presentation.
    This function extracts just the custom show definitions (name, id, slide_numbers).
    
    Args:
        pptx_path: Path to the PPTX file
    
    Returns:
        Tuple of (custom_shows dict, set of slide numbers used by custom shows)
        
        custom_shows structure:
        {
            0: {"name": "romans6.3", "id": 0, "slide_numbers": [8]},
            1: {"name": "revelation20.6", "id": 1, "slide_numbers": [16, 17]}
        }
    """
    custom_shows = {}
    custom_show_slide_nums = set()  # Track which slide numbers are used by custom shows
    
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            pres_xml = zf.read('ppt/presentation.xml').decode('utf-8')
            root = ET.fromstring(pres_xml)
            
            # Parse relationships to map rId to slide number
            rels_xml = zf.read('ppt/_rels/presentation.xml.rels').decode('utf-8')
            rels_root = ET.fromstring(rels_xml)
            rid_to_slide_num = {}
            for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                r_id = rel.get('Id')
                target = rel.get('Target')
                if target and 'slide' in target.lower() and not 'layout' in target.lower() and not 'master' in target.lower():
                    # Extract slide number from target (e.g., "slides/slide6.xml" -> 6)
                    import re
                    match = re.search(r'slide(\d+)\.xml', target)
                    if match:
                        rid_to_slide_num[r_id] = int(match.group(1))
            
            # Parse custom shows - just extract metadata (name, id, slide_numbers)
            for custShow in root.findall('.//p:custShow', NAMESPACES):
                show_name = custShow.get('name', '')
                show_id = custShow.get('id')
                
                if show_id:
                    slide_numbers = []
                    for sld in custShow.findall('.//p:sld', NAMESPACES):
                        r_id = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        if r_id and r_id in rid_to_slide_num:
                            slide_num = rid_to_slide_num[r_id]
                            slide_numbers.append(slide_num)
                            custom_show_slide_nums.add(slide_num)
                    
                    custom_shows[int(show_id)] = {
                        'name': show_name,
                        'id': int(show_id),
                        'slide_numbers': slide_numbers
                    }
        except Exception as e:
            print(f"Error parsing custom shows: {e}")
    
    return custom_shows, custom_show_slide_nums


def process_linked_slide(slide_num, prs, zf, slide_width, theme_colors, file_path, custom_shows=None):
    """Process a single linked slide (from custom show or hlinksldjump) and return its content.
    
    Args:
        slide_num: The 1-based slide number to process
        prs: The Presentation object
        zf: Open ZipFile object for the PPTX
        slide_width: Slide width in pixels
        theme_colors: Theme color map
        file_path: Path to the PPTX file (for image extraction)
        custom_shows: Dict of custom shows for name lookup
    
    Returns:
        Dict with slide_number, animation_sequence, and optionally static_content
    """
    slide_file = f'ppt/slides/slide{slide_num}.xml'
    slide_xml = zf.read(slide_file).decode('utf-8')
    
    # Get the actual slide object for shape access
    slide_obj = prs.slides[slide_num - 1]
    
    # Build shape ID to pptx shape mapping
    pptx_shapes_by_id = {}
    for z_idx, shape, group_id in enumerate_shapes_recursive(slide_obj.shapes):
        pptx_shapes_by_id[str(shape.shape_id)] = (shape, z_idx, group_id)
    
    # Build z_index mapping for all shapes (including pictures)
    shape_z_indices = {shape_id: data[1] for shape_id, data in pptx_shapes_by_id.items()}
    
    # Get animation entries
    animation_entries = parse_animation_sequence(slide_xml)
    
    # Get all shapes
    shapes = parse_shapes_from_slide(slide_xml)
    
    # Get all pictures from this slide
    rid_to_image = parse_slide_relationships(file_path, slide_num)
    pictures = parse_pictures_from_slide(slide_xml, rid_to_image, slide_width, shape_z_indices)
    
    # Get slide hyperlink relationships (for resolving hlinksldjump)
    rid_to_target_slide = parse_slide_links_from_relationships(file_path, slide_num)
    
    # Build animation sequence
    animation_sequence = []
    sequence_num = 1
    
    for anim_entry in animation_entries:
        shape_id = anim_entry['shape_id']
        
        # Check if this shape_id is a group
        child_ids = get_group_child_ids(slide_xml, shape_id)
        
        if child_ids:
            # Group animation - add all children
            for child_id in child_ids:
                if child_id in shapes:
                    shape = shapes[child_id]
                    entry = {
                        'sequence': sequence_num,
                        'shape_name': shape['name'],
                        'timing': anim_entry['timing']
                    }
                    has_text = bool(shape['text'])
                    if has_text:
                        entry['text'] = shape['text']
                    
                    if 'delay' in anim_entry and anim_entry['delay'] > 0:
                        entry['delay'] = anim_entry['delay']
                    
                    # Add visual data
                    if child_id in pptx_shapes_by_id:
                        pptx_shape, z_idx, grp_id = pptx_shapes_by_id[child_id]
                        visual = extract_shape_visual_data(pptx_shape, z_idx, slide_xml, child_id, slide_width, theme_colors)
                        if visual:
                            for key, value in visual.items():
                                if key == 'font' and not has_text:
                                    continue
                                entry[key] = value
                        entry['group_id'] = shape_id
                    
                    if shape['hyperlink']:
                        hyperlink = shape['hyperlink'].copy()
                        if hyperlink['type'] == 'slide':
                            r_id = hyperlink.get('r_id')
                            if r_id and r_id in rid_to_target_slide:
                                hyperlink['slide_number'] = rid_to_target_slide[r_id]
                                del hyperlink['r_id']
                        elif hyperlink['type'] == 'customshow' and custom_shows:
                            show_id = hyperlink.get('id')
                            if show_id is not None and show_id in custom_shows:
                                hyperlink['name'] = custom_shows[show_id]['name']
                        entry['hyperlink'] = hyperlink

                    _sanitize_scripture_block_font(entry, shape.get('is_scripture', False))
                    
                    animation_sequence.append(entry)
            sequence_num += 1
        elif shape_id in shapes:
            # Regular individual shape
            shape = shapes[shape_id]
            entry = {
                'sequence': sequence_num,
                'shape_name': shape['name'],
                'timing': anim_entry['timing']
            }
            has_text = bool(shape['text'])
            if has_text:
                entry['text'] = shape['text']
            
            if 'delay' in anim_entry and anim_entry['delay'] > 0:
                entry['delay'] = anim_entry['delay']
            
            # Add visual data
            if shape_id in pptx_shapes_by_id:
                pptx_shape, z_idx, grp_id = pptx_shapes_by_id[shape_id]
                visual = extract_shape_visual_data(pptx_shape, z_idx, slide_xml, shape_id, slide_width, theme_colors)
                if visual:
                    for key, value in visual.items():
                        if key == 'font' and not has_text:
                            continue
                        entry[key] = value
                if grp_id:
                    entry['group_id'] = grp_id
            
            if shape['hyperlink']:
                hyperlink = shape['hyperlink'].copy()
                if hyperlink['type'] == 'slide':
                    r_id = hyperlink.get('r_id')
                    if r_id and r_id in rid_to_target_slide:
                        hyperlink['slide_number'] = rid_to_target_slide[r_id]
                        del hyperlink['r_id']
                elif hyperlink['type'] == 'customshow' and custom_shows:
                    show_id = hyperlink.get('id')
                    if show_id is not None and show_id in custom_shows:
                        hyperlink['name'] = custom_shows[show_id]['name']
                entry['hyperlink'] = hyperlink

            _sanitize_scripture_block_font(entry, shape.get('is_scripture', False))
            
            animation_sequence.append(entry)
            sequence_num += 1
        elif shape_id in pictures:
            # Animated picture
            pic = pictures[shape_id]
            entry = {
                'sequence': sequence_num,
                'shape_name': pic['name'],
                'shape_type': 'picture',
                'image': pic['image'],
                'timing': anim_entry['timing']
            }
            if pic.get('description'):
                entry['description'] = pic['description']
            if 'delay' in anim_entry and anim_entry['delay'] > 0:
                entry['delay'] = anim_entry['delay']
            if pic.get('z_index') is not None:
                entry['z_index'] = pic['z_index']
            if pic.get('layout'):
                entry['layout'] = pic['layout']
            if pic.get('line'):
                entry['line'] = pic['line']
            animation_sequence.append(entry)
            sequence_num += 1
    
    # Get static content (non-animated shapes)
    static_shapes = []
    animated_ids = set(e['shape_id'] for e in animation_entries)
    for anim_entry in animation_entries:
        child_ids = get_group_child_ids(slide_xml, anim_entry['shape_id'])
        if child_ids:
            animated_ids.update(child_ids)
    
    for shape_id, shape in shapes.items():
        if shape_id not in animated_ids:
            static_entry = {
                'shape_name': shape['name'],
                'static': True
            }
            has_text = bool(shape['text'])
            if has_text:
                static_entry['text'] = shape['text']
            
            # Add visual data
            if shape_id in pptx_shapes_by_id:
                pptx_shape, z_idx, grp_id = pptx_shapes_by_id[shape_id]
                visual = extract_shape_visual_data(pptx_shape, z_idx, slide_xml, shape_id, slide_width, theme_colors)
                if visual:
                    for key, value in visual.items():
                        if key == 'font' and not has_text:
                            continue
                        static_entry[key] = value
                if grp_id:
                    static_entry['group_id'] = grp_id
            
            if shape['hyperlink']:
                hyperlink = shape['hyperlink'].copy()
                if hyperlink['type'] == 'slide':
                    r_id = hyperlink.get('r_id')
                    if r_id and r_id in rid_to_target_slide:
                        hyperlink['slide_number'] = rid_to_target_slide[r_id]
                        del hyperlink['r_id']
                elif hyperlink['type'] == 'customshow' and custom_shows:
                    show_id = hyperlink.get('id')
                    if show_id is not None and show_id in custom_shows:
                        hyperlink['name'] = custom_shows[show_id]['name']
                static_entry['hyperlink'] = hyperlink

            _sanitize_scripture_block_font(static_entry, shape.get('is_scripture', False))
            static_shapes.append(static_entry)
    
    # Add static pictures
    for pic_id, pic in pictures.items():
        if pic_id not in animated_ids:
            static_entry = {
                'shape_name': pic['name'],
                'shape_type': 'picture',
                'image': pic['image'],
                'static': True
            }
            if pic.get('description'):
                static_entry['description'] = pic['description']
            if pic.get('z_index') is not None:
                static_entry['z_index'] = pic['z_index']
            if pic.get('layout'):
                static_entry['layout'] = pic['layout']
            if pic.get('line'):
                static_entry['line'] = pic['line']
            static_shapes.append(static_entry)
    
    slide_content = {
        'slide_number': slide_num,
        'animation_sequence': animation_sequence
    }
    if static_shapes:
        slide_content['static_content'] = static_shapes
    
    return slide_content


def save_presentation_structure(prs, file_path):
    """Save a simplified representation focusing on animation order and hyperlinks."""
    
    # Extract theme colors from the PowerPoint file
    theme_colors = extract_theme_colors_from_pptx(file_path)
    
    if theme_colors:
        print(f"Extracted {len(theme_colors)} theme colors from presentation")
    
    # Parse custom shows - just metadata (name, id, slide_numbers)
    custom_shows, custom_show_slide_nums = parse_custom_shows(file_path)

    # Build linked-slide set from hidden slides only.
    # Non-hidden slides are always treated as top-level main slides.
    hidden_slide_nums = get_hidden_slides(file_path, len(prs.slides))
    linked_slide_nums = hidden_slide_nums

    if hidden_slide_nums:
        print(f"Found {len(hidden_slide_nums)} hidden slides: {sorted(hidden_slide_nums)}")
    
    # Calculate slide dimensions in pixels (from EMU)
    slide_width = emu_to_px(prs.slide_width)
    slide_height = emu_to_px(prs.slide_height)
    
    # Build linked_slides from hidden slide numbers only.
    linked_slides = {}
    with zipfile.ZipFile(file_path, 'r') as zf:
        for linked_slide_num in sorted(linked_slide_nums):
            try:
                slide_content = process_linked_slide(linked_slide_num, prs, zf, slide_width, theme_colors, file_path, custom_shows)
                linked_slides[linked_slide_num] = slide_content
            except Exception as e:
                print(f"Warning: Could not process linked slide {linked_slide_num}: {e}")
                linked_slides[linked_slide_num] = {
                    'slide_number': linked_slide_num,
                    'error': str(e)
                }
    
    # Calculate scale factor for coordinate conversion
    scale_factor = TARGET_CANVAS_WIDTH / slide_width if slide_width else 1.0
    
    # Create output folder for images (same name as JSON file)
    script_dir = Path(os.path.dirname(os.path.abspath(__file__)))
    extracted_dir = script_dir / 'extracted'
    extracted_dir.mkdir(exist_ok=True)
    
    output_stem = Path(file_path).stem
    images_folder = extracted_dir / output_stem
    
    # Extract all images from PPTX to the images folder
    image_map = extract_images_from_pptx(file_path, images_folder)
    
    # Calculate main presentation slide count (excluding linked slides)
    main_slide_count = len(prs.slides) - len(linked_slide_nums)
    
    # Store path with ~ instead of absolute home dir (safe for git, unambiguous)
    abs_path = Path(file_path).resolve()
    home = Path.home()
    try:
        display_path = '~/' + str(abs_path.relative_to(home))
    except ValueError:
        display_path = str(abs_path)

    presentation_data = {
        "file_path": display_path,
        "file_name": Path(file_path).name,
        "total_slides": main_slide_count,
        "total_custom_shows": len(custom_shows),
        "total_linked_slides": len(linked_slides),
        "source_dimensions": {
            "width": slide_width,
            "height": slide_height
        },
        "target_canvas": {
            "width": TARGET_CANVAS_WIDTH,
            "height": TARGET_CANVAS_HEIGHT
        },
        "scale_factor": round(scale_factor, 3),
        "images_folder": output_stem,
        "custom_shows": custom_shows,
        "linked_slides": linked_slides,
        "slides": []
    }
    
    # Build a mapping of shape ID to pptx shape object for visual data
    with zipfile.ZipFile(file_path, 'r') as zf:
        for slide_num, slide in enumerate(prs.slides, 1):
            # Skip linked slides (they go into linked_slides, not main presentation)
            if slide_num in linked_slide_nums:
                continue
                
            slide_file = f'ppt/slides/slide{slide_num}.xml'
            
            # Build shape ID to pptx shape mapping for this slide (including grouped shapes)
            pptx_shapes_by_id = {}
            for z_idx, shape, group_id in enumerate_shapes_recursive(slide.shapes):
                pptx_shapes_by_id[str(shape.shape_id)] = (shape, z_idx, group_id)
            
            # Build z_index mapping for all shapes (including pictures)
            shape_z_indices = {shape_id: data[1] for shape_id, data in pptx_shapes_by_id.items()}
            
            try:
                slide_xml = zf.read(slide_file).decode('utf-8')
                
                # Get animation entries with timing info
                animation_entries = parse_animation_sequence(slide_xml)
                
                # Get all shapes
                shapes = parse_shapes_from_slide(slide_xml)
                
                # Get all pictures from this slide
                rid_to_image = parse_slide_relationships(file_path, slide_num)
                pictures = parse_pictures_from_slide(slide_xml, rid_to_image, slide_width, shape_z_indices)
                
                # Get slide link relationships for resolving hlinksldjump
                rid_to_target_slide = parse_slide_links_from_relationships(file_path, slide_num)
                
                # Build ordered animation list
                animation_sequence = []
                sequence_num = 1
                
                for anim_entry in animation_entries:
                    shape_id = anim_entry['shape_id']
                    
                    # Check if this shape_id is a group
                    child_ids = get_group_child_ids(slide_xml, shape_id)
                    
                    if child_ids:
                        # This is a group - add all child shapes with the same sequence/timing
                        for child_id in child_ids:
                            if child_id in shapes:
                                shape = shapes[child_id]
                                entry = {
                                    'sequence': sequence_num,
                                    'shape_name': shape['name']
                                }
                                has_text = bool(shape['text'])
                                if has_text:
                                    entry['text'] = shape['text']
                                
                                # Add timing info (all children get same timing)
                                entry['timing'] = anim_entry['timing']
                                if 'delay' in anim_entry and anim_entry['delay'] > 0:
                                    entry['delay'] = anim_entry['delay']
                                
                                # Add visual data if available (with auto-scaling)
                                if child_id in pptx_shapes_by_id:
                                    pptx_shape, z_idx, group_id = pptx_shapes_by_id[child_id]
                                    visual = extract_shape_visual_data(pptx_shape, z_idx, slide_xml, child_id, slide_width, theme_colors)
                                    if visual:
                                        for key, value in visual.items():
                                            # Skip font for shapes without text
                                            if key == 'font' and not has_text:
                                                continue
                                            entry[key] = value
                                    # Mark that this is part of an animated group
                                    entry['group_id'] = shape_id
                                
                                # Add hyperlink info if present (just reference, no content)
                                if shape['hyperlink']:
                                    hyperlink = shape['hyperlink'].copy()
                                    if hyperlink['type'] == 'slide':
                                        # Resolve r_id to target slide number
                                        r_id = hyperlink.get('r_id')
                                        if r_id and r_id in rid_to_target_slide:
                                            hyperlink['slide_number'] = rid_to_target_slide[r_id]
                                            del hyperlink['r_id']
                                    elif hyperlink['type'] == 'customshow':
                                        show_id = hyperlink.get('id')
                                        if show_id is not None and show_id in custom_shows:
                                            hyperlink['name'] = custom_shows[show_id]['name']
                                    entry['hyperlink'] = hyperlink

                                _sanitize_scripture_block_font(entry, shape.get('is_scripture', False))
                                
                                animation_sequence.append(entry)
                        sequence_num += 1
                    elif shape_id in shapes:
                        # Regular individual shape
                        shape = shapes[shape_id]
                        # Include all animated shapes (text or not - could be rectangles, decorative shapes)
                        entry = {
                            'sequence': sequence_num,
                            'shape_name': shape['name']
                        }
                        has_text = bool(shape['text'])
                        if has_text:
                            entry['text'] = shape['text']
                        
                        # Add timing info (click, with, after)
                        entry['timing'] = anim_entry['timing']
                        if 'delay' in anim_entry and anim_entry['delay'] > 0:
                            entry['delay'] = anim_entry['delay']
                        
                        # Add visual data if available (with auto-scaling)
                        if shape_id in pptx_shapes_by_id:
                            pptx_shape, z_idx, group_id = pptx_shapes_by_id[shape_id]
                            visual = extract_shape_visual_data(pptx_shape, z_idx, slide_xml, shape_id, slide_width, theme_colors)
                            if visual:
                                for key, value in visual.items():
                                    # Skip font for shapes without text
                                    if key == 'font' and not has_text:
                                        continue
                                    entry[key] = value
                            # Add group_id for debugging if shape is in a group
                            if group_id:
                                entry['group_id'] = group_id
                        
                        # Add hyperlink info if present (just reference, no content)
                        if shape['hyperlink']:
                            hyperlink = shape['hyperlink'].copy()
                            if hyperlink['type'] == 'slide':
                                # Resolve r_id to target slide number
                                r_id = hyperlink.get('r_id')
                                if r_id and r_id in rid_to_target_slide:
                                    hyperlink['slide_number'] = rid_to_target_slide[r_id]
                                    del hyperlink['r_id']
                            elif hyperlink['type'] == 'customshow':
                                show_id = hyperlink.get('id')
                                if show_id is not None and show_id in custom_shows:
                                    hyperlink['name'] = custom_shows[show_id]['name']
                            entry['hyperlink'] = hyperlink

                        _sanitize_scripture_block_font(entry, shape.get('is_scripture', False))
                        
                        animation_sequence.append(entry)
                        sequence_num += 1
                    elif shape_id in pictures:
                        # Animated picture
                        pic = pictures[shape_id]
                        entry = {
                            'sequence': sequence_num,
                            'shape_name': pic['name'],
                            'shape_type': 'picture',
                            'image': pic['image'],
                            'timing': anim_entry['timing']
                        }
                        if pic.get('description'):
                            entry['description'] = pic['description']
                        if 'delay' in anim_entry and anim_entry['delay'] > 0:
                            entry['delay'] = anim_entry['delay']
                        if pic.get('z_index') is not None:
                            entry['z_index'] = pic['z_index']
                        if pic.get('layout'):
                            entry['layout'] = pic['layout']
                        if pic.get('line'):
                            entry['line'] = pic['line']
                        
                        animation_sequence.append(entry)
                        sequence_num += 1
                
                # Also get shapes that might not be animated (static content)
                # Build a set of ALL animated shape IDs (including those in animated groups)
                static_shapes = []
                animated_ids = set(e['shape_id'] for e in animation_entries)
                # Also add child IDs from animated groups
                for anim_entry in animation_entries:
                    child_ids = get_group_child_ids(slide_xml, anim_entry['shape_id'])
                    if child_ids:
                        animated_ids.update(child_ids)
                
                for shape_id, shape in shapes.items():
                    if shape_id not in animated_ids:
                        # Include all shapes - text, connectors, or decorative rectangles
                        # Skip if it's truly empty (no text, no visual importance)
                        static_entry = {
                            'shape_name': shape['name'],
                            'static': True
                        }
                        has_text = bool(shape['text'])
                        if has_text:
                            static_entry['text'] = shape['text']
                        
                        # Add visual data if available (with auto-scaling)
                        if shape_id in pptx_shapes_by_id:
                            pptx_shape, z_idx, group_id = pptx_shapes_by_id[shape_id]
                            visual = extract_shape_visual_data(pptx_shape, z_idx, slide_xml, shape_id, slide_width, theme_colors)
                            if visual:
                                for key, value in visual.items():
                                    # Skip font for shapes without text
                                    if key == 'font' and not has_text:
                                        continue
                                    static_entry[key] = value
                            # Add group_id for debugging if shape is in a group
                            if group_id:
                                static_entry['group_id'] = group_id
                        else:
                            # Fallback: Extract visual data from XML for shapes not in pptx enumeration
                            # (e.g., shapes in groups that python-pptx doesn't expose)
                            try:
                                root = ET.fromstring(slide_xml)
                                # Find shape element by ID
                                for sp in root.findall('.//p:sp', NAMESPACES):
                                    cNvPr = sp.find('.//p:cNvPr', NAMESPACES)
                                    if cNvPr is not None and cNvPr.get('id') == shape_id:
                                        # Extract visual data from XML
                                        visual = extract_visual_data_from_xml(sp, 0, slide_width)
                                        if visual:
                                            for key, value in visual.items():
                                                # Skip font for shapes without text
                                                if key == 'font' and not has_text:
                                                    continue
                                                static_entry[key] = value
                                        # Check if shape is in a group by looking for parent grpSp
                                        parent = sp
                                        for _ in range(5):  # Check up to 5 levels
                                            parent = parent.find('..')
                                            if parent is None:
                                                break
                                            if parent.tag.endswith('grpSp'):
                                                grp_cNvPr = parent.find('.//p:cNvPr', NAMESPACES)
                                                if grp_cNvPr is not None:
                                                    static_entry['group_id'] = grp_cNvPr.get('id')
                                                break
                                        break
                            except:
                                pass
                        
                        if shape['hyperlink']:
                            hyperlink = shape['hyperlink'].copy()
                            if hyperlink['type'] == 'slide':
                                # Resolve r_id to target slide number
                                r_id = hyperlink.get('r_id')
                                if r_id and r_id in rid_to_target_slide:
                                    hyperlink['slide_number'] = rid_to_target_slide[r_id]
                                    del hyperlink['r_id']
                            elif hyperlink['type'] == 'customshow':
                                show_id = hyperlink.get('id')
                                if show_id is not None and show_id in custom_shows:
                                    hyperlink['name'] = custom_shows[show_id]['name']
                            static_entry['hyperlink'] = hyperlink

                        _sanitize_scripture_block_font(static_entry, shape.get('is_scripture', False))
                        static_shapes.append(static_entry)
                
                # Add static pictures (non-animated)
                for pic_id, pic in pictures.items():
                    if pic_id not in animated_ids:
                        static_entry = {
                            'shape_name': pic['name'],
                            'shape_type': 'picture',
                            'image': pic['image'],
                            'static': True
                        }
                        if pic.get('description'):
                            static_entry['description'] = pic['description']
                        if pic.get('z_index') is not None:
                            static_entry['z_index'] = pic['z_index']
                        if pic.get('layout'):
                            static_entry['layout'] = pic['layout']
                        if pic.get('line'):
                            static_entry['line'] = pic['line']
                        static_shapes.append(static_entry)
                
                slide_info = {
                    'slide_number': slide_num,
                    'animation_sequence': animation_sequence,
                }
                
                if static_shapes:
                    slide_info['static_content'] = static_shapes
                
                presentation_data['slides'].append(slide_info)
                
            except Exception as e:
                presentation_data['slides'].append({
                    'slide_number': slide_num,
                    'error': str(e)
                })
    
    # Save JSON file inside the same subfolder as images
    script_dir = Path(os.path.dirname(os.path.abspath(__file__)))
    extracted_dir = script_dir / 'extracted'
    output_folder = extracted_dir / output_stem
    output_folder.mkdir(parents=True, exist_ok=True)
    
    output_filename = output_stem + '.json'
    output_path = output_folder / output_filename
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(presentation_data, f, indent=2, ensure_ascii=False)
    
    print(f"Presentation structure saved to: {output_path}")
    return output_path


def get_pptx_file():
    script_directory = os.path.dirname(os.path.abspath(__file__))
   
    #expect hsu-pptx or pptx folder to be in the same directory as this script
    path = Path(script_directory).parent / 'hsu-pptx'
    if not path.is_dir():
        path = Path(script_directory).parent / 'pptx'    
    if not path.is_dir():
        print(f"Error. Files not found in: {Path(script_directory).parent}\n" 
              f"Add a folder named 'pptx' to the same directory as this script and add .pptx files to it.")
        exit()
    # Use glob to filter and sort .pptx files
    file_list = sorted(glob.glob(os.path.join(path, '*.pptx')))

    # Print the list of files to the console
    if file_list:
        print(f"Extract from: ${path}")
        for index, file in enumerate(file_list):
            print(f"{index + 1}. {Path(file).name}")

   # Ask the user to select a file
    while True:
        try:
            selection = int(input("Enter the number of the file you want to extract text from (0 to exit): "))
            
            # Check if the selection is valid
            if 0 <= selection <= len(file_list):
                if selection == 0:
                    print("Exiting...")
                    exit()
                else:
                    selected_file = file_list[selection - 1]
                    print(f"Selected: {selected_file}")
                    return selected_file

            else:
                print("Invalid selection. Please enter a valid number.")
        except ValueError:
            print("Invalid input. Please enter a valid number.")

def main():
    import sys
    
    # Accept file path as command-line argument, or fall back to interactive selection
    if len(sys.argv) > 1:
        file_name = sys.argv[1]
        if not os.path.exists(file_name):
            print(f"Error: File not found: {file_name}")
            exit(1)
        if not file_name.endswith('.pptx'):
            print(f"Error: File must be a .pptx file: {file_name}")
            exit(1)
        print(f"Processing: {file_name}")
    else:
        file_name = get_pptx_file()
    
    # Load the presentation
    prs = Presentation(file_name)
    
    save_presentation_structure(prs, file_name)

if __name__ == "__main__":
    main()